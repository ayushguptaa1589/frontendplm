"""
PLM System - Comprehensive Presentation Generator
Generates a professional PowerPoint covering all features, security, roles, dashboards, etc.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──────────────────────────────────────────
DARK_BG      = RGBColor(0x0F, 0x17, 0x2A)   # Deep navy
ACCENT_BLUE  = RGBColor(0x00, 0x7B, 0xFF)   # Primary blue
ACCENT_CYAN  = RGBColor(0x00, 0xD4, 0xFF)   # Bright cyan
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY     = RGBColor(0x88, 0x99, 0xAA)
CARD_BG      = RGBColor(0x1A, 0x25, 0x3C)   # Card background
GREEN        = RGBColor(0x28, 0xA7, 0x45)
ORANGE       = RGBColor(0xFF, 0x9F, 0x43)
RED          = RGBColor(0xDC, 0x35, 0x45)
PURPLE       = RGBColor(0x6F, 0x42, 0xC1)
YELLOW       = RGBColor(0xFF, 0xD9, 0x3D)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper Functions ───────────────────────────────────────

def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_fill(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = 0
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_paragraph(text_frame, text, font_size=16, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(4), font_name='Segoe UI'):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p

def add_bullet_list(slide, left, top, width, height, items, font_size=15, color=LIGHT_GRAY, bullet_color=ACCENT_CYAN, title=None, title_size=20):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(title_size)
        p.font.color.rgb = ACCENT_CYAN
        p.font.bold = True
        p.font.name = 'Segoe UI'
        p.space_after = Pt(12)
        first = False
    else:
        first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = f"  ●  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Segoe UI'
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    return txBox

def add_card(slide, left, top, width, height, title, items, icon="", card_color=CARD_BG, title_color=ACCENT_CYAN):
    shape = add_shape_fill(slide, left, top, width, height, card_color)
    shape.shadow.inherit = False
    # Rounded corners via shape
    inner_left = left + Inches(0.25)
    inner_top = top + Inches(0.15)
    inner_w = width - Inches(0.5)
    
    txBox = slide.shapes.add_textbox(inner_left, inner_top, inner_w, height - Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{icon}  {title}" if icon else title
    p.font.size = Pt(17)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = 'Segoe UI'
    p.space_after = Pt(8)
    
    for item in items:
        p2 = tf.add_paragraph()
        p2.text = f"  ▸  {item}"
        p2.font.size = Pt(13)
        p2.font.color.rgb = LIGHT_GRAY
        p2.font.name = 'Segoe UI'
        p2.space_before = Pt(3)
        p2.space_after = Pt(3)
    return shape

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = add_shape_fill(slide, left, top, width, Inches(0.04), color)
    return shape

def slide_header(slide, title, subtitle=None):
    set_slide_bg(slide)
    add_accent_line(slide, Inches(0.6), Inches(0.5), Inches(3.5), ACCENT_BLUE)
    add_text_box(slide, Inches(0.6), Inches(0.6), Inches(10), Inches(0.7), title, font_size=36, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(1.2), Inches(10), Inches(0.5), subtitle, font_size=18, color=MID_GRAY)

# ══════════════════════════════════════════════════════════════
#   SLIDE 1: Title Slide
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

# Accent bar
add_shape_fill(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT_BLUE)

# Title block
add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
             "PLM", font_size=72, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(0.8),
             "Product Lifecycle Management System", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(1.5), Inches(3.45), Inches(5), ACCENT_CYAN)
add_text_box(slide, Inches(1.5), Inches(3.7), Inches(10), Inches(0.6),
             "Complete System Overview  |  Features  |  Security  |  Architecture", font_size=18, color=MID_GRAY)
add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.5),
             "Built with Node.js  ·  Express.js  ·  SQLite  ·  Vanilla JS", font_size=16, color=LIGHT_GRAY)
add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.4),
             "Pravaig  |  February 2026", font_size=14, color=MID_GRAY)

# ══════════════════════════════════════════════════════════════
#   SLIDE 2: Table of Contents
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Table of Contents")

sections = [
    ("01", "System Overview & Architecture"),
    ("02", "Technology Stack"),
    ("03", "User Roles & Permissions"),
    ("04", "Login & Signup System"),
    ("05", "Admin Dashboard"),
    ("06", "Designer Dashboard"),
    ("07", "Approver Dashboard"),
    ("08", "Vault: Parts & Assemblies"),
    ("09", "Version Control & BOM"),
    ("10", "Engineering Change Orders (ECOs)"),
    ("11", "Approval Workflow"),
    ("12", "Security System"),
    ("13", "Notifications & Search"),
    ("14", "Data Export & Analytics"),
    ("15", "UI/UX Features"),
    ("16", "Database Schema"),
    ("17", "API Endpoints Summary"),
    ("18", "Deployment & Infrastructure"),
]

col1 = sections[:9]
col2 = sections[9:]

for i, (num, title) in enumerate(col1):
    y = Inches(1.9) + Inches(i * 0.55)
    add_text_box(slide, Inches(1.0), y, Inches(0.6), Inches(0.4), num, font_size=16, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(1.6), y, Inches(4.5), Inches(0.4), title, font_size=16, color=LIGHT_GRAY)

for i, (num, title) in enumerate(col2):
    y = Inches(1.9) + Inches(i * 0.55)
    add_text_box(slide, Inches(7.0), y, Inches(0.6), Inches(0.4), num, font_size=16, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(7.6), y, Inches(4.5), Inches(0.4), title, font_size=16, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════════
#   SLIDE 3: System Overview & Architecture
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "System Overview & Architecture", "End-to-end product lifecycle management for engineering teams")

# Left description
add_bullet_list(slide, Inches(0.6), Inches(1.9), Inches(5.5), Inches(4.5), [
    "Web-based PLM system for managing parts, assemblies & products",
    "3-tier role hierarchy: Admin → Approver → Designer",
    "Complete approval workflow for designs & engineering changes",
    "Version-controlled vault for all engineering data",
    "Built as a Single Page Application (SPA) architecture",
    "RESTful API backend with JWT token authentication",
    "SQLite embedded database — zero external dependencies",
    "Real-time notifications and activity logging",
], font_size=15)

# Right: Architecture box
add_card(slide, Inches(6.8), Inches(1.9), Inches(5.8), Inches(4.8),
         "Architecture Layers", [
             "Frontend: HTML5, CSS3, Vanilla JavaScript",
             "API Layer: Express.js REST endpoints (74 routes)",
             "Auth Layer: JWT tokens + bcrypt password hashing",
             "Security: Helmet, Rate Limiting, CORS, Compression",
             "Database: SQLite3 with 16+ tables & indexes",
             "File Storage: Multer disk storage (50MB limit)",
             "Middleware: Token verification, Role guards, Validation",
         ], icon="🏗️")

# ══════════════════════════════════════════════════════════════
#   SLIDE 4: Technology Stack
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Technology Stack", "Modern, lightweight yet powerful tooling")

cards_data = [
    ("🖥️ Backend", ACCENT_BLUE, [
        "Node.js v24.x runtime",
        "Express.js web framework",
        "SQLite3 embedded database",
        "JWT (jsonwebtoken) auth",
        "bcryptjs password hashing",
    ]),
    ("🎨 Frontend", ACCENT_CYAN, [
        "HTML5 semantic markup",
        "CSS3 with CSS variables",
        "Vanilla JavaScript (ES6+)",
        "Responsive dark/light theme",
        "Toast notification system",
    ]),
    ("🔒 Security", GREEN, [
        "Helmet HTTP headers",
        "express-rate-limit",
        "CORS cross-origin policy",
        "Compression (gzip)",
        "Input validation & enums",
    ]),
    ("📦 Infrastructure", ORANGE, [
        "Multer file uploads",
        "Disk-based file storage",
        "Auto DB schema migration",
        "Graceful shutdown handling",
        "Health check endpoint",
    ]),
]

for i, (title, title_color, items) in enumerate(cards_data):
    x = Inches(0.5) + Inches(i * 3.15)
    add_card(slide, x, Inches(2.0), Inches(2.9), Inches(4.0), title, items, title_color=title_color)

# ══════════════════════════════════════════════════════════════
#   SLIDE 5: User Roles & Permissions
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "User Roles & Permissions", "Three-tier role-based access control system")

# Admin card
add_card(slide, Inches(0.5), Inches(2.0), Inches(3.8), Inches(5.0),
         "👑  ADMIN", [
             "Full system control & configuration",
             "Approve/reject new user signups",
             "Create, edit, delete any project",
             "Manage all users (activate/deactivate)",
             "Access all vault items (parts & assemblies)",
             "Create & manage Engineering Change Orders",
             "View analytics dashboard & reports",
             "Export data to CSV",
             "Bulk delete/freeze operations",
             "Review all approval requests",
             "View complete activity history",
         ], title_color=ACCENT_BLUE)

# Approver card
add_card(slide, Inches(4.7), Inches(2.0), Inches(3.8), Inches(5.0),
         "✅  APPROVER", [
             "Review & approve/reject submissions",
             "Approve edit access requests",
             "Approve release requests",
             "Freeze part & assembly versions",
             "View vault items (parts & assemblies)",
             "Impact analysis for parts",
             "BOM (Bill of Materials) access",
             "Bulk freeze operations",
             "Provide feedback on submissions",
             "View submission history",
         ], title_color=GREEN)

# Designer card
add_card(slide, Inches(8.9), Inches(2.0), Inches(3.8), Inches(5.0),
         "🎨  DESIGNER", [
             "Create & manage own designs",
             "Submit designs for approval",
             "Create parts & assemblies in vault",
             "Create new part/assembly versions",
             "Request edit access to parts",
             "Request part/assembly releases",
             "Track submission status",
             "View project progress",
             "Rollback part versions",
             "View vault & BOM data",
         ], title_color=ORANGE)

# ══════════════════════════════════════════════════════════════
#   SLIDE 6: Login & Signup System
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Login & Signup System", "Secure authentication with admin-gated registration")

add_card(slide, Inches(0.5), Inches(2.0), Inches(5.8), Inches(4.8),
         "🔐  Authentication Flow", [
             "Role-based login: select Admin / Designer / Approver",
             "Username + Password + Role validated on server",
             "JWT token generated (24-hour expiry)",
             "Token stored in localStorage for session persistence",
             "Auto-redirect to role-specific dashboard",
             "Inactive accounts blocked at login",
             "Rate limited: 20 attempts per 15 minutes",
         ], icon="", title_color=ACCENT_BLUE)

add_card(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.8),
         "📝  Signup & Account Management", [
             "New users register with username, email, password, role",
             "Accounts created as INACTIVE by default",
             "Admin must approve each new signup",
             "Admin can reject signups with feedback",
             "Password requirements: 8+ chars, uppercase, lowercase, number, special",
             "Change password (requires current password)",
             "Change email (requires password verification)",
             "Forgot password with reset code system",
         ], icon="", title_color=GREEN)

# ══════════════════════════════════════════════════════════════
#   SLIDE 7: Admin Dashboard
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Admin Dashboard", "Central command center for system administration")

modules = [
    ("✅ Approvals", ["Approve/reject pending user signups", "View signup details & role", "Activate or reject accounts"]),
    ("📊 Projects", ["Create, edit, delete projects", "Track progress with percentage", "Set deadlines & assign managers"]),
    ("👥 Users", ["View all registered users", "Filter by role and status", "Edit user roles, delete users"]),
    ("🗂️ Vault", ["Browse all parts & assemblies", "Create/edit/delete vault items", "View versions, BOM, impact"]),
    ("📝 Requests", ["Review edit access requests", "Review release requests", "Approve/reject with feedback"]),
    ("📜 History", ["Full activity log for all users", "Filter by user, action type, date", "Timestamped audit trail (IST)"]),
    ("📈 Reports", ["7 live analytics stat cards", "Parts by lifecycle bar chart", "Export parts/assemblies CSV"]),
    ("🔄 Change Orders", ["Create ECOs with priority levels", "Track ECO status workflow", "Comment threads on ECOs"]),
]

for i, (title, items) in enumerate(modules):
    col = i % 4
    row = i // 4
    x = Inches(0.4) + Inches(col * 3.15)
    y = Inches(2.0) + Inches(row * 2.6)
    add_card(slide, x, y, Inches(3.0), Inches(2.4), title, items, title_color=ACCENT_CYAN)

# ══════════════════════════════════════════════════════════════
#   SLIDE 8: Designer Dashboard
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Designer Dashboard", "Design creation, submission, and vault management")

modules = [
    ("📊 My Projects", ["View assigned projects", "Track project progress", "Update progress percentage"]),
    ("🎨 My Designs", ["Create new designs", "Submit designs for review", "Track submission status"]),
    ("🗂️ Vault", ["Create parts with metadata", "Create assemblies with BOM", "Version control (create/rollback)"]),
    ("📝 Requests", ["Request edit access to parts", "Request part/assembly release", "Track request status"]),
    ("📜 History", ["View personal activity log", "Track all design submissions", "Timestamped action history"]),
    ("⚙️ Profile", ["Change password securely", "Update email with verification", "View role & account info"]),
]

for i, (title, items) in enumerate(modules):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(col * 4.1)
    y = Inches(2.0) + Inches(row * 2.7)
    add_card(slide, x, y, Inches(3.8), Inches(2.4), title, items, title_color=ORANGE)

# ══════════════════════════════════════════════════════════════
#   SLIDE 9: Approver Dashboard
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Approver Dashboard", "Review, approve, and govern engineering outputs")

modules = [
    ("✅ Pending Approvals", ["View all pending submissions", "Approve or reject with feedback", "Decision recorded with timestamp"]),
    ("📊 Projects", ["View all active projects", "Monitor overall project health", "Track deadlines & progress"]),
    ("🗂️ Vault", ["Browse parts & assemblies", "Freeze versions for release", "View impact analysis & BOM"]),
    ("📝 Requests", ["Review edit access requests", "Review release requests", "Grant/deny with reasoning"]),
    ("📜 History", ["Complete approval audit trail", "Filter decisions by status", "Timestamped records (IST)"]),
    ("⚙️ Profile", ["Secure password change", "Email update with verification", "Account settings management"]),
]

for i, (title, items) in enumerate(modules):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(col * 4.1)
    y = Inches(2.0) + Inches(row * 2.7)
    add_card(slide, x, y, Inches(3.8), Inches(2.4), title, items, title_color=GREEN)

# ══════════════════════════════════════════════════════════════
#   SLIDE 10: Vault - Parts & Assemblies
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Vault: Parts & Assemblies", "Engineering data vault with full lifecycle management")

add_card(slide, Inches(0.5), Inches(2.0), Inches(5.8), Inches(5.0),
         "🔩  Parts Management", [
             "Create parts with: code, name, description",
             "Metadata: material, vendor, criticality, tags",
             "Lifecycle states: Draft → Active → Released → Obsolete",
             "Version control with version labels",
             "Freeze versions to lock changes",
             "Rollback to previous versions",
             "Impact analysis: see which assemblies use a part",
             "Edit permissions system per part",
             "Edit/release request workflow",
             "Bulk delete & bulk freeze (admin/approver)",
         ], icon="", title_color=ACCENT_BLUE)

add_card(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(5.0),
         "🏗️  Assembly Management", [
             "Create assemblies with: code, name, description",
             "Metadata: criticality, lifecycle state, tags",
             "Bill of Materials (BOM) for each version",
             "BOM maps assembly versions → part versions",
             "Version control per assembly",
             "Freeze assembly versions",
             "Lifecycle states match parts lifecycle",
             "Import/export BOM data as CSV",
             "View all parts contained in an assembly",
             "Criticality levels: Normal, Low, High, Critical",
         ], icon="", title_color=GREEN)

# ══════════════════════════════════════════════════════════════
#   SLIDE 11: Version Control & BOM
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Version Control & BOM", "Engineering-grade version management and bill of materials")

add_card(slide, Inches(0.5), Inches(2.0), Inches(3.8), Inches(5.0),
         "📌  Version Control", [
             "Each part/assembly has multiple versions",
             'Status: "Working" or "Frozen"',
             "Freeze locks a version permanently",
             "Frozen by user + timestamp recorded",
             "Version labels (e.g., v1.0, v2.0)",
             "Rollback to any previous version",
             "Change notes per version",
             "Compare two versions side by side",
             "Storage path & working path tracking",
         ], icon="", title_color=ACCENT_CYAN)

add_card(slide, Inches(4.7), Inches(2.0), Inches(3.8), Inches(5.0),
         "📋  Bill of Materials", [
             "BOM ties assembly versions to part versions",
             "View complete part list per assembly version",
             "Part details in BOM: code, name, material, vendor",
             "Version status visible in BOM view",
             "Export BOM to CSV for any assembly version",
             "Supports multi-level product structure",
             "Critical for manufacturing & procurement",
         ], icon="", title_color=ORANGE)

add_card(slide, Inches(8.9), Inches(2.0), Inches(3.8), Inches(5.0),
         "🔍  Impact Analysis", [
             "See which assemblies reference a part",
             "Critical for change assessment",
             "Shows assembly name, code, version",
             "Prevents accidental deletion of used parts",
             "Bulk delete checks assembly references",
             "Essential for ECO planning",
             "Helps approvers assess change risk",
         ], icon="", title_color=PURPLE)

# ══════════════════════════════════════════════════════════════
#   SLIDE 12: Engineering Change Orders (ECOs)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Engineering Change Orders (ECOs)", "Formal change management workflow")

add_card(slide, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.3),
         "📝  ECO Creation", [
             "Title, description, reason for change",
             "Priority: Low, Medium, High, Critical",
             "Specify affected parts & assemblies",
             "Auto-generated unique ECO number (ECO-timestamp)",
         ], icon="", title_color=ACCENT_BLUE)

add_card(slide, Inches(0.5), Inches(4.5), Inches(5.8), Inches(2.6),
         "🔄  ECO Workflow Status", [
             "Draft → Submitted → In Review → Approved / Rejected → Implemented",
             "Status updates trigger notifications to admins & approvers",
             "reviewer_id tracks who is reviewing",
             "decided_at records decision timestamp",
             "Implementation notes for final closure",
         ], icon="", title_color=GREEN)

add_card(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(5.1),
         "💬  ECO Collaboration", [
             "Comment threads on each ECO",
             "Any authenticated user can comment",
             "Comments show username + timestamp",
             "Discussion visible in ECO detail modal",
             "Search ECOs by number, title, or description",
             "Filter by status and priority",
             "Delete ECOs (admin only)",
             "Notifications sent on status transitions",
         ], icon="", title_color=ORANGE)

# ══════════════════════════════════════════════════════════════
#   SLIDE 13: Approval Workflow
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Approval Workflow", "Multi-stage approval process for quality governance")

# Flow diagram using text
add_text_box(slide, Inches(0.6), Inches(2.0), Inches(12), Inches(0.5),
             "DESIGN SUBMISSION FLOW", font_size=20, color=ACCENT_CYAN, bold=True)

flow_steps = [
    ("Designer", "Creates design &\nsubmits for review", ORANGE),
    ("→", "", MID_GRAY),
    ("Approver", "Reviews design,\napproves/rejects", GREEN),
    ("→", "", MID_GRAY),
    ("Admin", "Final oversight\n& project management", ACCENT_BLUE),
]

for i, (title, desc, color) in enumerate(flow_steps):
    x = Inches(0.5) + Inches(i * 2.5)
    if title == "→":
        add_text_box(slide, x, Inches(3.0), Inches(1.0), Inches(0.5), "→", font_size=40, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    else:
        shape = add_shape_fill(slide, x, Inches(2.7), Inches(2.2), Inches(1.5), CARD_BG)
        add_text_box(slide, x + Inches(0.15), Inches(2.8), Inches(1.9), Inches(0.4), title, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.15), Inches(3.2), Inches(1.9), Inches(0.8), desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Additional workflows
add_card(slide, Inches(0.5), Inches(4.6), Inches(3.8), Inches(2.5),
         "📝  Edit Request Flow", [
             "Designer requests edit access",
             "Admin/Approver reviews request",
             "Approve → grant edit permission",
             "Reject → designer cannot edit",
         ], title_color=ORANGE)

add_card(slide, Inches(4.7), Inches(4.6), Inches(3.8), Inches(2.5),
         "🚀  Release Request Flow", [
             "Designer requests part/assembly release",
             "Approver reviews & validates",
             "Approve → lifecycle → Released",
             "Reject → back to designer",
         ], title_color=GREEN)

add_card(slide, Inches(8.9), Inches(4.6), Inches(3.8), Inches(2.5),
         "👤  User Signup Flow", [
             "New user registers (inactive)",
             "Admin sees pending signup",
             "Approve → account activated",
             "Reject → account remains inactive",
         ], title_color=ACCENT_BLUE)

# ══════════════════════════════════════════════════════════════
#   SLIDE 14: Security System
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Security System", "Multi-layered security architecture")

add_card(slide, Inches(0.5), Inches(2.0), Inches(3.8), Inches(5.0),
         "🔐  Authentication", [
             "JWT tokens with 24-hour expiry",
             "bcrypt password hashing (salt rounds: 10)",
             "Token verification on every API call",
             "Role embedded in JWT payload",
             "Automatic session expiry & redirect",
             "Inactive account rejection at login",
             "Password strength enforcement",
             "   - 8+ characters minimum",
             "   - Uppercase + lowercase required",
             "   - Number + special char required",
         ], title_color=RED)

add_card(slide, Inches(4.7), Inches(2.0), Inches(3.8), Inches(5.0),
         "🛡️  API Security", [
             "Helmet.js HTTP security headers",
             "  - X-Content-Type-Options",
             "  - X-Frame-Options",
             "  - Strict-Transport-Security",
             "Rate limiting on all API routes",
             "  - General: 500 req / 15 min",
             "  - Auth: 20 req / 15 min",
             "CORS cross-origin protection",
             "JSON payload size limit (10MB)",
             "Gzip compression for responses",
             "Global error handler (no stack leaks)",
         ], title_color=ORANGE)

add_card(slide, Inches(8.9), Inches(2.0), Inches(3.8), Inches(5.0),
         "🔑  Access Control", [
             "Role-based route guards (middleware)",
             "requireRole() enforces per-route access",
             "Enum validation on all inputs",
             "  - Roles: admin, designer, approver",
             "  - Criticality: normal, low, high, critical",
             "  - Lifecycle: draft, active, released, obsolete",
             "  - Priority: low, medium, high, critical",
             "  - ECO Status: 6 valid states",
             "File upload type whitelist",
             "  - PDF, STEP, STL, DWG, DXF, etc.",
             "50MB max file upload size",
         ], title_color=GREEN)

# ══════════════════════════════════════════════════════════════
#   SLIDE 15: Notifications & Search
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Notifications & Global Search", "Stay informed and find anything instantly")

add_card(slide, Inches(0.5), Inches(2.0), Inches(5.8), Inches(5.0),
         "🔔  Notification System", [
             "Real-time notification bell in top bar",
             "Unread count badge with auto-refresh (30s)",
             "Notification dropdown with recent 20 items",
             "Mark all as read with one click",
             "Notification types: info, success, warning, error",
             "Auto-generated on ECO status changes",
             "Notifies admins & approvers on new ECOs",
             "Toast notifications for all user actions",
             "  - Success (green), Error (red), Warning (orange), Info (blue)",
             "  - Auto-dismiss after configurable duration",
             "  - Close button for manual dismiss",
         ], icon="", title_color=ACCENT_CYAN)

add_card(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(5.0),
         "🔍  Global Search", [
             "Search bar available on all dashboards",
             "Searches across 5 entity types simultaneously:",
             "  - Parts (by code, name, material, tags)",
             "  - Assemblies (by code, name, tags)",
             "  - Projects (by PLM ID, name)",
             "  - Users (by username, email)",
             "  - ECOs (by number, title, description)",
             "Debounced input (300ms) for performance",
             "Categorized dropdown results",
             "Minimum 2 characters to trigger search",
             "Results limited to 10 per category",
         ], icon="", title_color=ACCENT_BLUE)

# ══════════════════════════════════════════════════════════════
#   SLIDE 16: Data Export & Analytics
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Data Export & Analytics", "Business intelligence and data portability")

add_card(slide, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.3),
         "📊  Analytics Dashboard (Admin)", [
             "7 live stat cards: Projects, Users, Parts, Assemblies, ECOs, Edit Requests, Release Requests",
             "Parts by lifecycle state bar chart visualization",
             "All data fetched in real-time from database",
         ], icon="", title_color=ACCENT_BLUE)

add_card(slide, Inches(0.5), Inches(4.5), Inches(5.8), Inches(2.5),
         "📥  CSV Export", [
             "Export all parts with metadata to CSV",
             "Export all assemblies with metadata to CSV",
             "Export BOM for any specific assembly version",
             "One-click download buttons in Reports section",
         ], icon="", title_color=GREEN)

add_card(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(5.0),
         "📜  Activity History & Audit Trail", [
             "Every action logged with: user, action, type, timestamp",
             "Activity types: create, update, delete, login, approval",
             "Timestamps converted to India Standard Time (IST)",
             "Filterable by user, action type, date range",
             "Complete audit trail for compliance",
             "Auto-logged: logins, signups, approvals, CRUD operations",
             "Supports forensic analysis and accountability",
             "Print-friendly view with @media print styles",
         ], icon="", title_color=ORANGE)

# ══════════════════════════════════════════════════════════════
#   SLIDE 17: UI/UX Features
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "UI/UX Features", "Polished, modern interface with attention to detail")

cards = [
    ("🌗  Dark / Light Theme", [
        "Toggle switch in top bar",
        "CSS variables for easy theming",
        "Preference saved in localStorage",
        "Persists across sessions",
    ]),
    ("🔔  Toast Notifications", [
        "4 types: success, error, warning, info",
        "Animated slide-in/out",
        "Auto-dismiss with timer",
        "Manual close button",
    ]),
    ("📱  Responsive Layout", [
        "Sidebar navigation",
        "Content sections toggled by menu",
        "Modal dialogs for forms",
        "Overflow scroll for large data",
    ]),
    ("🖨️  Print-Friendly", [
        "Print stylesheet included",
        "Hides sidebar, nav, modals",
        "White background for printing",
        "Page margins & borders auto-set",
    ]),
    ("⏳  Loading Spinners", [
        "Visual feedback during API calls",
        "Animated CSS spinner",
        "Shown in content containers",
        "Improves perceived performance",
    ]),
    ("🔗  Hierarchy Flow", [
        "Visual role hierarchy display",
        "Designer → Approver → Admin",
        "Shows on all dashboards",
        "Highlights current user's role",
    ]),
]

for i, (title, items) in enumerate(cards):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(col * 4.15)
    y = Inches(2.0) + Inches(row * 2.7)
    add_card(slide, x, y, Inches(3.85), Inches(2.4), title, items, title_color=ACCENT_CYAN)

# ══════════════════════════════════════════════════════════════
#   SLIDE 18: Database Schema
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Database Schema", "16 tables with indexes for performance • SQLite3")

tables_left = [
    ("users", "id, username, email, password, role, is_active, approved_by, created_at, reset_code"),
    ("projects", "id, plm_id, name, owner_id, status, progress, deadline, manager, created_at"),
    ("parts", "id, part_code, name, description, material, vendor, criticality, lifecycle_state, tags, owner_id"),
    ("part_versions", "id, part_id, version_label, status, storage_path, working_path, change_notes, frozen_by/at"),
    ("assemblies", "id, assembly_code, name, description, criticality, lifecycle_state, tags, owner_id"),
    ("assembly_versions", "id, assembly_id, version_label, status, storage_path, change_notes, frozen_by/at"),
    ("assembly_parts", "id, assembly_version_id, part_version_id (BOM mapping)"),
    ("part_permissions", "id, part_id, user_id, can_edit (per-part access control)"),
]

tables_right = [
    ("submissions", "id, project_id, designer_id, submission_type, status, file_path, comments"),
    ("approvals", "id, submission_id, approver_id, decision, feedback, decision_date"),
    ("tasks", "id, project_id, designer_id, title, description, priority, due_date, completed"),
    ("activity_logs", "id, user_id, username, action, action_type, details, timestamp"),
    ("notifications", "id, user_id, title, message, type, is_read, link, created_at"),
    ("eco_orders", "id, eco_number, title, description, reason, priority, status, requester_id, reviewer_id"),
    ("comments", "id, entity_type, entity_id, user_id, username, message, created_at"),
    ("attachments", "id, entity_type, entity_id, filename, original_name, file_path, file_size, mime_type"),
]

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(6.0), Inches(5.2))
tf = txBox.text_frame
tf.word_wrap = True
for i, (table, cols) in enumerate(tables_left):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = f"  {table}"
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_CYAN
    p.font.bold = True
    p.font.name = 'Segoe UI Semibold'
    p.space_before = Pt(6)
    
    p2 = tf.add_paragraph()
    p2.text = f"     {cols}"
    p2.font.size = Pt(11)
    p2.font.color.rgb = MID_GRAY
    p2.font.name = 'Segoe UI'
    p2.space_after = Pt(4)

txBox2 = slide.shapes.add_textbox(Inches(6.8), Inches(1.9), Inches(6.0), Inches(5.2))
tf2 = txBox2.text_frame
tf2.word_wrap = True
for i, (table, cols) in enumerate(tables_right):
    if i == 0:
        p = tf2.paragraphs[0]
    else:
        p = tf2.add_paragraph()
    p.text = f"  {table}"
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_CYAN
    p.font.bold = True
    p.font.name = 'Segoe UI Semibold'
    p.space_before = Pt(6)
    
    p2 = tf2.add_paragraph()
    p2.text = f"     {cols}"
    p2.font.size = Pt(11)
    p2.font.color.rgb = MID_GRAY
    p2.font.name = 'Segoe UI'
    p2.space_after = Pt(4)

# Indexes note
add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
             "19 database indexes optimized for: user lookups, activity logs, parts lifecycle, versions, submissions, approvals, ECOs, notifications, comments, attachments",
             font_size=12, color=MID_GRAY)

# ══════════════════════════════════════════════════════════════
#   SLIDE 19: API Endpoints Summary
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "API Endpoints Summary", "74 RESTful API routes • All JWT-protected (except auth & health)")

api_groups = [
    ("Authentication (3)", "POST /login, /signup, /forgot-password, /reset-password"),
    ("Parts CRUD (9)", "POST/GET/PUT/DELETE /parts, versions, freeze, rollback, impact, permissions"),
    ("Assemblies CRUD (8)", "POST/GET/PUT/DELETE /assemblies, versions, freeze, BOM"),
    ("Edit Requests (4)", "POST create, GET list, POST approve, POST reject"),
    ("Release Requests (4)", "POST create, GET list, POST approve, POST reject"),
    ("Projects (4)", "GET list, POST create, PUT update, DELETE remove"),
    ("Tasks (2)", "GET list, POST create"),
    ("Submissions (3)", "GET list, POST create, DELETE remove"),
    ("Approvals (2)", "GET pending, POST decide"),
    ("Users (6)", "GET list, GET pending, POST approve/reject, PUT edit, DELETE remove"),
    ("Account (3)", "POST change-password, POST update-email, GET current-user"),
    ("Activity (1)", "GET /activity-history with filters"),
    ("Notifications (3)", "GET list, GET unread-count, POST mark-read"),
    ("ECOs (4)", "GET list, POST create, PUT update, DELETE (admin)"),
    ("Comments (2)", "GET by entity, POST to entity"),
    ("Files (3)", "POST upload, GET attachments, GET download"),
    ("Export (3)", "GET /export/parts, /assemblies, /bom/:id/:version"),
    ("Search (1)", "GET /search?q= (cross-entity)"),
    ("Analytics (1)", "GET /analytics/dashboard"),
    ("Bulk Ops (2)", "POST bulk-delete, POST bulk-freeze"),
    ("Versions (1)", "GET /versions/compare?v1=&v2="),
    ("Health (1)", "GET /health (public, DB check)"),
]

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(5.8), Inches(5.2))
tf = txBox.text_frame
tf.word_wrap = True
first = True
for group, desc in api_groups[:11]:
    if first:
        p = tf.paragraphs[0]
        first = False
    else:
        p = tf.add_paragraph()
    p.text = f"  {group}"
    p.font.size = Pt(13)
    p.font.color.rgb = ACCENT_CYAN
    p.font.bold = True
    p.font.name = 'Segoe UI'
    p.space_before = Pt(4)
    p2 = tf.add_paragraph()
    p2.text = f"     {desc}"
    p2.font.size = Pt(11)
    p2.font.color.rgb = MID_GRAY
    p2.font.name = 'Segoe UI'
    p2.space_after = Pt(2)

txBox2 = slide.shapes.add_textbox(Inches(6.8), Inches(1.9), Inches(5.8), Inches(5.2))
tf2 = txBox2.text_frame
tf2.word_wrap = True
first = True
for group, desc in api_groups[11:]:
    if first:
        p = tf2.paragraphs[0]
        first = False
    else:
        p = tf2.add_paragraph()
    p.text = f"  {group}"
    p.font.size = Pt(13)
    p.font.color.rgb = ACCENT_CYAN
    p.font.bold = True
    p.font.name = 'Segoe UI'
    p.space_before = Pt(4)
    p2 = tf2.add_paragraph()
    p2.text = f"     {desc}"
    p2.font.size = Pt(11)
    p2.font.color.rgb = MID_GRAY
    p2.font.name = 'Segoe UI'
    p2.space_after = Pt(2)

# ══════════════════════════════════════════════════════════════
#   SLIDE 20: Deployment & Infrastructure
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Deployment & Infrastructure", "Easy to deploy, maintain, and scale")

add_card(slide, Inches(0.5), Inches(2.0), Inches(3.8), Inches(4.6),
         "🚀  Quick Start", [
             "npm install (one command)",
             "node server.js (starts on :5000)",
             "Or use start.bat on Windows",
             "Auto-creates database on first run",
             "Auto-creates uploads/ directory",
             "Default users seeded automatically",
             "Zero external database dependencies",
         ], title_color=GREEN)

add_card(slide, Inches(4.7), Inches(2.0), Inches(3.8), Inches(4.6),
         "📁  Project Structure", [
             "server.js — Backend (2,667 lines)",
             "api.js — Shared API helpers",
             "index.html — Login/signup page",
             "admin-dashboard.html + .js",
             "designer-dashboard.html + .js",
             "approver-dashboard.html + .js",
             "dashboard.css + style.css",
             ".gitignore — Security protection",
         ], title_color=ACCENT_CYAN)

add_card(slide, Inches(8.9), Inches(2.0), Inches(3.8), Inches(4.6),
         "⚙️  Production Features", [
             "Graceful shutdown (SIGTERM/SIGINT)",
             "10-second forced exit timeout",
             "Database connection close on shutdown",
             "Health check endpoint for monitoring",
             "gzip compression for all responses",
             ".gitignore protects: .env, .db, uploads/",
             "Environment variable support (.env)",
         ], title_color=ORANGE)

# ══════════════════════════════════════════════════════════════
#   SLIDE 21: Summary / Thank You
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_shape_fill(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT_BLUE)

add_text_box(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(0.8),
             "PLM System — Complete Summary", font_size=40, color=WHITE, bold=True)
add_accent_line(slide, Inches(1.5), Inches(2.1), Inches(5), ACCENT_CYAN)

stats = [
    ("74", "API Routes"),
    ("16", "Database Tables"),
    ("3", "User Roles"),
    ("3", "Dashboards"),
    ("6", "ECO States"),
    ("4", "Lifecycle States"),
    ("19", "DB Indexes"),
    ("2,667", "Lines of Backend"),
]

for i, (num, label) in enumerate(stats):
    col = i % 4
    row = i // 4
    x = Inches(1.5) + Inches(col * 2.8)
    y = Inches(2.6) + Inches(row * 1.6)
    add_text_box(slide, x, y, Inches(2.0), Inches(0.6), num, font_size=42, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.6), Inches(2.0), Inches(0.4), label, font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.5),
             "A complete Product Lifecycle Management system — from login to production-ready deployment.",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.4),
             "Pravaig  •  Built with Node.js, Express, SQLite  •  February 2026",
             font_size=14, color=MID_GRAY, alignment=PP_ALIGN.LEFT)

# ── Save ──────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'PLM_System_Overview.pptx')
prs.save(output_path)
print(f"✅ Presentation saved: {output_path}")
print(f"   Slides: {len(prs.slides)}")
